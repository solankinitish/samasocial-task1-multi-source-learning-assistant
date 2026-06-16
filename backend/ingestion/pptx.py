from pptx import Presentation
from backend.schemas.models import Chunk

def ingest_pptx(file_path: str, source_label: str) -> list[Chunk]:
    try:
        chunks = []
        prs = Presentation(file_path)
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            slide_text += text + " "
            slide_text = slide_text.strip()
            if slide_text:
                chunks.append(Chunk(
                    text=slide_text,
                    source_type="pptx",
                    source_label=source_label,
                    location=f"slide {slide_num}"
                ))
        return chunks
    except Exception as e:
        raise ValueError(f"Failed to process PPTX {source_label}: {str(e)}")
